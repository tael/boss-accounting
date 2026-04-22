from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── 색상 팔레트 ──────────────────────────────────────
BRAND   = RGBColor(0x10, 0x99, 0x48)   # 초록 (앱 컬러)
DARK    = RGBColor(0x1F, 0x29, 0x37)
MID     = RGBColor(0x6B, 0x72, 0x80)
LIGHT   = RGBColor(0xF9, 0xFA, 0xFB)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT  = RGBColor(0x25, 0x63, 0xEB)   # 파랑
WARN    = RGBColor(0xDC, 0x26, 0x26)   # 빨강

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank = prs.slide_layouts[6]  # 완전 빈 레이아웃

# ── 헬퍼 ─────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=None, line=None):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background() if line is None else None
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = wrap
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def header_bar(slide, title, subtitle=""):
    add_rect(slide, 0, 0, prs.slide_width, Inches(1.1), fill=BRAND)
    add_text(slide, title, Inches(0.4), Inches(0.12), Inches(10), Inches(0.55),
             size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(0.65), Inches(12), Inches(0.38),
                 size=14, color=RGBColor(0xBB, 0xF7, 0xD0))

def footer(slide, page_num):
    add_rect(slide, 0, prs.slide_height - Inches(0.35), prs.slide_width, Inches(0.35),
             fill=RGBColor(0xF3, 0xF4, 0xF6))
    add_text(slide, "사장님 회계 도우미 — Gemini 코드리뷰 개선 보고서",
             Inches(0.4), prs.slide_height - Inches(0.32), Inches(10), Inches(0.3),
             size=9, color=MID)
    add_text(slide, str(page_num),
             Inches(12.8), prs.slide_height - Inches(0.32), Inches(0.4), Inches(0.3),
             size=9, color=MID, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════
# 슬라이드 1: 표지
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, fill=DARK)
add_rect(slide, 0, 0, Inches(0.35), prs.slide_height, fill=BRAND)

add_text(slide, "사장님 회계 도우미",
         Inches(0.8), Inches(1.5), Inches(12), Inches(0.9),
         size=44, bold=True, color=WHITE)
add_text(slide, "Gemini 코드리뷰 기반 클린코드 개선 보고서",
         Inches(0.8), Inches(2.5), Inches(11), Inches(0.5),
         size=22, color=RGBColor(0xBB, 0xF7, 0xD0))
add_text(slide, "React 19 · TypeScript · Zustand · Tailwind CSS 4",
         Inches(0.8), Inches(3.1), Inches(10), Inches(0.4),
         size=14, color=MID)

# 통계 박스 3개
stats = [
    ("8 Rounds", "Gemini 리뷰 사이클"),
    ("37 Stories", "개선 항목"),
    ("50 → 96", "테스트 건수"),
]
for i, (num, label) in enumerate(stats):
    x = Inches(0.8 + i * 3.5)
    add_rect(slide, x, Inches(4.3), Inches(3.1), Inches(1.5),
             fill=RGBColor(0x37, 0x41, 0x51))
    add_text(slide, num, x + Inches(0.15), Inches(4.45), Inches(2.8), Inches(0.55),
             size=30, bold=True, color=BRAND, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + Inches(0.15), Inches(5.05), Inches(2.8), Inches(0.4),
             size=12, color=MID, align=PP_ALIGN.CENTER)

add_text(slide, "2026.04.22", Inches(0.8), Inches(6.5), Inches(4), Inches(0.4),
         size=11, color=MID)

# ══════════════════════════════════════════════════════
# 슬라이드 2: 배경 & 목표
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
header_bar(slide, "개선 배경 & 목표", "왜 Gemini 코드리뷰를 도입했는가")
footer(slide, 2)

add_rect(slide, Inches(0.3), Inches(1.3), prs.slide_width - Inches(0.6), Inches(5.7),
         fill=LIGHT)

cols = [
    ("🔍 문제 인식", ACCENT, [
        "뷰/로직/스토어 경계가 불명확",
        "컴포넌트 내부 비즈니스 로직 혼재",
        "중복 코드 & 불일치 패턴 산재",
        "잠재적 런타임 버그 존재",
    ]),
    ("🎯 개선 원칙", BRAND, [
        "클린코드 — 관심사 분리 철저",
        "베스트 프랙티스 준수",
        "수정 범위 최소화",
        "정책·밸런스 변경 용이성",
    ]),
    ("🔄 방법론", WARN, [
        "Gemini CLI v0.27.3 코드리뷰",
        "스토리별 AC 수립 후 구현",
        "매 라운드 tsc + 테스트 검증",
        "커밋 단위 점진적 개선",
    ]),
]

for i, (title, color, items) in enumerate(cols):
    x = Inches(0.5 + i * 4.1)
    add_rect(slide, x, Inches(1.4), Inches(3.9), Inches(0.45), fill=color)
    add_text(slide, title, x + Inches(0.1), Inches(1.43), Inches(3.7), Inches(0.4),
             size=14, bold=True, color=WHITE)
    for j, item in enumerate(items):
        add_text(slide, f"• {item}", x + Inches(0.1), Inches(2.0 + j * 0.7),
                 Inches(3.7), Inches(0.65), size=12, color=DARK)

# ══════════════════════════════════════════════════════
# 슬라이드 3: 라운드 전체 개요 (타임라인)
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
header_bar(slide, "8라운드 개선 타임라인", "Gemini 리뷰 → 계획 → 구현 → 검증 반복")
footer(slide, 3)

rounds = [
    ("R1", "스토어 캡슐화\n중복 제거", BRAND),
    ("R2", "로직 추출\n(VAT/소득세)", ACCENT),
    ("R3", "타입 정리\n래퍼 제거", RGBColor(0x70, 0x55, 0xE3)),
    ("R4", "구글드라이브\n안전성 강화", RGBColor(0xF5, 0x9E, 0x0B)),
    ("R5", "접근성\n(A11y) 개선", WARN),
    ("R6", "런타임 버그\n수정 5건", RGBColor(0xEC, 0x48, 0x99)),
    ("R7", "계산 버그\n수정 3건", RGBColor(0x06, 0xB6, 0xD4)),
    ("R8", "Export/Import\nYAxis/scope", RGBColor(0x84, 0xCC, 0x16)),
]

y_base = Inches(1.4)
for i, (label, desc, color) in enumerate(rounds):
    x = Inches(0.35 + i * 1.55)
    add_rect(slide, x, y_base, Inches(1.35), Inches(0.5), fill=color)
    add_text(slide, label, x, y_base + Inches(0.04), Inches(1.35), Inches(0.42),
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # 화살표 연결선 (텍스트로)
    if i < len(rounds) - 1:
        add_text(slide, "→", x + Inches(1.36), y_base + Inches(0.08), Inches(0.18), Inches(0.35),
                 size=16, color=MID, align=PP_ALIGN.CENTER)

# 각 라운드 설명
for i, (label, desc, color) in enumerate(rounds):
    x = Inches(0.35 + i * 1.55)
    add_text(slide, desc, x, y_base + Inches(0.6), Inches(1.35), Inches(0.8),
             size=9.5, color=DARK, align=PP_ALIGN.CENTER)

# 결과 요약 박스
add_rect(slide, Inches(0.3), Inches(3.1), prs.slide_width - Inches(0.6), Inches(3.9),
         fill=LIGHT)
add_text(slide, "라운드별 주요 성과",
         Inches(0.5), Inches(3.2), Inches(12), Inches(0.4),
         size=14, bold=True, color=DARK)

results = [
    ("R1-2", "스토어 캡슐화 완성\nimportTransactions, bulkAddTransactions 추가\nformatAmount·12개월 생성 로직 중복 제거"),
    ("R3-4", "관심사 분리 완성\ngetUpcoming → taxSchedule.ts 이동\ngetFiltered store action → filterTransactions util"),
    ("R5-6", "안전성 강화\nToast aria-live, aria-pressed 추가\nTAX_DEDUCTION_CONSTANTS 상수화"),
    ("R7-8", "버그 수정 완료\nBreakEvenCalc variableRatio 수정\nVatCalculator 런타임 크래시 수정"),
]

for i, (title, detail) in enumerate(results):
    x = Inches(0.4 + i * 3.1)
    add_rect(slide, x, Inches(3.65), Inches(2.9), Inches(3.1),
             fill=WHITE)
    add_text(slide, title, x + Inches(0.1), Inches(3.72), Inches(2.7), Inches(0.35),
             size=13, bold=True, color=BRAND)
    add_text(slide, detail, x + Inches(0.1), Inches(4.1), Inches(2.7), Inches(2.5),
             size=10, color=DARK)

# ══════════════════════════════════════════════════════
# 슬라이드 4: 스토어 & 아키텍처 개선 (R1-2)
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
header_bar(slide, "Round 1-2: 아키텍처 개선", "스토어 캡슐화 · 중복 제거 · 로직 분리")
footer(slide, 4)

# Before / After
add_text(slide, "Before", Inches(0.3), Inches(1.3), Inches(5.9), Inches(0.4),
         size=16, bold=True, color=WARN)
add_rect(slide, Inches(0.3), Inches(1.7), Inches(5.9), Inches(5.4),
         fill=RGBColor(0xFF, 0xF1, 0xF2))

before_items = [
    "❌  DataPage.tsx: useTransactionStore.setState() 직접 호출",
    "❌  TransactionForm: today 상수 모듈 레벨 → stale date 버그",
    "❌  DashboardPage: formatAmount() 인라인 중복 함수",
    "❌  TrendChart + CashFlowChart: 동일한 12개월 생성 로직 2벌",
    "❌  VatCalculator 컴포넌트 내 분기 필터링 도메인 로직",
    "❌  IncomeTaxCalc useMemo 내 절세 계산 직접 구현",
]
for i, item in enumerate(before_items):
    add_text(slide, item, Inches(0.5), Inches(1.85 + i * 0.75), Inches(5.5), Inches(0.65),
             size=10.5, color=WARN)

add_text(slide, "After", Inches(6.8), Inches(1.3), Inches(6.1), Inches(0.4),
         size=16, bold=True, color=BRAND)
add_rect(slide, Inches(6.8), Inches(1.7), Inches(6.1), Inches(5.4),
         fill=RGBColor(0xF0, 0xFD, 0xF4))

after_items = [
    "✅  importTransactions + bulkAddTransactions 액션 추가",
    "✅  useState(() => getTodayLocal()) — 컴포넌트 마운트 기준",
    "✅  formatKRW (utils/format.ts) 재사용",
    "✅  getLastNMonths(n) → financial.ts 공통 함수",
    "✅  filterByQuarter + summarizeForVat → financial.ts 순수 함수",
    "✅  calculateYellowUmbrellaSaving → tax.ts 순수 함수",
]
for i, item in enumerate(after_items):
    add_text(slide, item, Inches(7.0), Inches(1.85 + i * 0.75), Inches(5.7), Inches(0.65),
             size=10.5, color=BRAND)

# ══════════════════════════════════════════════════════
# 슬라이드 5: 타입 안전 & 래퍼 제거 (R3-4)
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
header_bar(slide, "Round 3-4: 타입 안전성 & 스토어 정리", "alias 제거 · 래퍼 메서드 제거 · googleDrive 안전화")
footer(slide, 5)

items = [
    (ACCENT, "IncomeTaxResult alias 필드 제거",
     "taxBase / tax / effectiveRate → taxableIncomeKRW / calculatedTaxKRW / appliedRatePct\n단일 canonical 필드명으로 통일"),
    (BRAND, "transactionStore 래퍼 메서드 제거",
     "getAll / getByDateRange / getByCategory 제거\n→ transactions 배열 + financial.ts 유틸 직접 사용 (Zustand 권장 패턴)"),
    (RGBColor(0x70, 0x55, 0xE3), "getFiltered store action → util 분리",
     "derived state 필터링은 store 밖으로\n→ filterTransactions(transactions, filter) 유틸 함수"),
    (RGBColor(0xF5, 0x9E, 0x0B), "googleDrive.ts setInterval 폴링 제거",
     "100ms × 100회 폴링 → script onload/onerror 이벤트 기반 Promise\n+ settled 플래그로 메모리 누수 방지"),
    (WARN, "TYPE_LABEL 강타입화",
     "Record<string, string> → Record<TaxType, string>\nTaxType 유니온 export → 컴파일 타임 키 검증"),
]

for i, (color, title, desc) in enumerate(items):
    col = i % 2
    row = i // 2
    x = Inches(0.3 + col * 6.5)
    y = Inches(1.3 + row * 2.2)
    w = Inches(6.2)
    h = Inches(2.0)
    add_rect(slide, x, y, w, h, fill=LIGHT)
    add_rect(slide, x, y, Inches(0.18), h, fill=color)
    add_text(slide, title, x + Inches(0.28), y + Inches(0.12), w - Inches(0.4), Inches(0.38),
             size=13, bold=True, color=DARK)
    add_text(slide, desc, x + Inches(0.28), y + Inches(0.55), w - Inches(0.4), Inches(1.3),
             size=10.5, color=MID)

# ══════════════════════════════════════════════════════
# 슬라이드 6: 안전성 & 접근성 (R5-6)
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
header_bar(slide, "Round 5-6: 안전성 & 접근성 개선", "A11y · 매직 넘버 상수화 · 런타임 버그 수정")
footer(slide, 6)

fixes = [
    ("🛡️ ErrorBoundary 추가", "런타임 크래시 시 빈 화면 대신 복구 UI 표시\nHashRouter를 ErrorBoundary로 감쌈"),
    ("♿ Toast aria-live", "에러: role=alert + aria-live=assertive\n일반: role=status + aria-live=polite"),
    ("♿ aria-pressed", "TransactionForm 매출/비용 토글 버튼\nOnboardingModal 선택 버튼에 aria-pressed 추가"),
    ("🔢 매직 넘버 상수화", "1_500_000 → PERSONAL_DEDUCTION_PER_DEPENDENT\n5_000_000 → YELLOW_UMBRELLA_MAX"),
    ("🐛 VatCalculator 크래시", "S-025에서 제거된 getFiltered 참조 잔존\n→ transactions 직접 사용으로 긴급 수정"),
    ("🐛 isMounted ref", "onSuccess 후 unmount된 컴포넌트에 setSubmitting(false) 호출\n→ isMounted.current 체크로 방지"),
    ("🐛 crypto.randomUUID", "비 HTTPS 환경(HTTP staging)에서 TypeError\n→ typeof 체크 + Math.random 폴백"),
    ("🐛 dependents 최솟값", "부양가족 0명 입력 시 세금 과다 계산\n→ Math.max(1, value)로 clamp"),
]

for i, (title, desc) in enumerate(fixes):
    col = i % 2
    row = i // 2
    x = Inches(0.25 + col * 6.5)
    y = Inches(1.3 + row * 1.45)
    add_rect(slide, x, y, Inches(6.2), Inches(1.35), fill=LIGHT)
    add_text(slide, title, x + Inches(0.15), y + Inches(0.1), Inches(5.9), Inches(0.38),
             size=12, bold=True, color=DARK)
    add_text(slide, desc, x + Inches(0.15), y + Inches(0.5), Inches(5.9), Inches(0.75),
             size=10, color=MID)

# ══════════════════════════════════════════════════════
# 슬라이드 7: 테스트 현황
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
header_bar(slide, "테스트 커버리지 현황", "96개 단위 테스트 — 5개 파일 전부 통과")
footer(slide, 7)

# 그래프 바
add_text(slide, "개선 전 (50건)", Inches(0.4), Inches(1.4), Inches(3), Inches(0.4),
         size=13, bold=True, color=MID)
add_rect(slide, Inches(0.4), Inches(1.85), Inches(5.2), Inches(0.7),
         fill=RGBColor(0xE5, 0xE7, 0xEB))
add_rect(slide, Inches(0.4), Inches(1.85), Inches(5.2), Inches(0.7),
         fill=MID)
add_text(slide, "50", Inches(5.75), Inches(1.9), Inches(0.8), Inches(0.6),
         size=20, bold=True, color=MID)

add_text(slide, "개선 후 (96건)", Inches(0.4), Inches(2.75), Inches(3), Inches(0.4),
         size=13, bold=True, color=BRAND)
add_rect(slide, Inches(0.4), Inches(3.2), Inches(10.0), Inches(0.7),
         fill=BRAND)
add_text(slide, "96", Inches(10.55), Inches(3.25), Inches(0.8), Inches(0.6),
         size=20, bold=True, color=BRAND)

# 파일별 현황
add_text(slide, "파일별 테스트 현황", Inches(0.4), Inches(4.15), Inches(12), Inches(0.4),
         size=14, bold=True, color=DARK)

test_files = [
    ("tax.test.ts", "17건", "세금 계산 (VAT, 종합소득세, BEP)", ACCENT),
    ("validation.test.ts", "16건", "입력값 검증 로직", RGBColor(0x70, 0x55, 0xE3)),
    ("financial.test.ts", "17건", "재무 계산 (손익, 현금흐름, 필터)", BRAND),
    ("format.test.ts", "21건", "포맷팅 유틸 (신규 추가)", RGBColor(0xF5, 0x9E, 0x0B)),
    ("exportImport.test.ts", "25건", "JSON 백업/복원 (신규 추가)", WARN),
]

for i, (fname, count, desc, color) in enumerate(test_files):
    x = Inches(0.3 + i * 2.5)
    add_rect(slide, x, Inches(4.65), Inches(2.3), Inches(2.5), fill=LIGHT)
    add_rect(slide, x, Inches(4.65), Inches(2.3), Inches(0.12), fill=color)
    add_text(slide, count, x, Inches(4.8), Inches(2.3), Inches(0.6),
             size=28, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, fname, x + Inches(0.1), Inches(5.48), Inches(2.1), Inches(0.35),
             size=9.5, bold=True, color=DARK)
    add_text(slide, desc, x + Inches(0.1), Inches(5.85), Inches(2.1), Inches(1.1),
             size=9, color=MID)

# ══════════════════════════════════════════════════════
# 슬라이드 8: 재무 계산 버그 수정 (R7)
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
header_bar(slide, "Round 7: 재무 계산 정확성 개선", "3개 버그 수정 — 실제 사용자 데이터에 영향")
footer(slide, 8)

bugs = [
    (
        "BreakEvenCalc variableRatio 오류",
        "변동비율 = 변동비 / (고정비+변동비)  ❌",
        "변동비율 = 변동비 / 매출  ✅",
        "BEP 공식 Fixed / (1-Ratio)는 매출 기준 비율 필요\n→ utils/financial.ts의 calculateBreakEven 직접 활용",
        WARN,
    ),
    (
        "sampleData.ts relativeDate UTC 버그",
        "toISOString().slice(0,10)  ❌",
        "로컬 getFullYear/getMonth/getDate  ✅",
        "KST 오전 0-9시 사용 시 예시 데이터가 하루 전 날짜로 생성\n→ getTodayLocal 패턴 적용",
        ACCENT,
    ),
    (
        "calculateCostStructure 0활동 월 스킵",
        "트랜잭션 있는 달만 targetMonths  ❌",
        "getLastNMonths(n) 기준 최근 N달  ✅",
        "활동 없는 달이 제외되어 고정비·매출 평균이 부풀려짐\n→ 전체 기간 포함, 빈 달은 0원으로 집계",
        BRAND,
    ),
]

for i, (title, before, after, reason, color) in enumerate(bugs):
    y = Inches(1.35 + i * 1.95)
    add_rect(slide, Inches(0.3), y, prs.slide_width - Inches(0.6), Inches(1.8), fill=LIGHT)
    add_rect(slide, Inches(0.3), y, Inches(0.18), Inches(1.8), fill=color)
    add_text(slide, title, Inches(0.6), y + Inches(0.1), Inches(12), Inches(0.35),
             size=13, bold=True, color=DARK)
    add_text(slide, f"Before: {before}", Inches(0.6), y + Inches(0.52),
             Inches(5.5), Inches(0.35), size=10.5, color=WARN)
    add_text(slide, f"After:  {after}", Inches(0.6), y + Inches(0.88),
             Inches(5.5), Inches(0.35), size=10.5, color=BRAND)
    add_text(slide, reason, Inches(6.5), y + Inches(0.3),
             Inches(6.4), Inches(1.2), size=10, color=MID)

# ══════════════════════════════════════════════════════
# 슬라이드 9: 개선 효과 요약
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
header_bar(slide, "개선 효과 요약", "코드 품질 지표 전후 비교")
footer(slide, 9)

metrics = [
    ("스토어 직접 setState", "3곳", "0곳", BRAND),
    ("컴포넌트 내 도메인 로직", "8개 함수", "0개", BRAND),
    ("중복 유틸 코드", "5건", "0건", BRAND),
    ("테스트 건수", "50건", "96건 (+92%)", BRAND),
    ("런타임 버그", "5건 잠재", "0건", BRAND),
    ("A11y 누락", "4건", "0건", BRAND),
    ("매직 넘버", "7개", "0개", BRAND),
    ("재무 계산 버그", "3건", "0건", BRAND),
]

add_rect(slide, Inches(0.3), Inches(1.25), prs.slide_width - Inches(0.6), Inches(0.45),
         fill=DARK)
add_text(slide, "지표", Inches(0.5), Inches(1.28), Inches(5.5), Inches(0.38),
         size=12, bold=True, color=WHITE)
add_text(slide, "개선 전", Inches(6.2), Inches(1.28), Inches(2.5), Inches(0.38),
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "개선 후", Inches(9.2), Inches(1.28), Inches(3.5), Inches(0.38),
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for i, (metric, before, after, color) in enumerate(metrics):
    bg = LIGHT if i % 2 == 0 else WHITE
    y = Inches(1.72 + i * 0.6)
    add_rect(slide, Inches(0.3), y, prs.slide_width - Inches(0.6), Inches(0.58), fill=bg)
    add_text(slide, metric, Inches(0.5), y + Inches(0.1), Inches(5.5), Inches(0.4),
             size=12, color=DARK)
    add_text(slide, before, Inches(6.2), y + Inches(0.1), Inches(2.5), Inches(0.4),
             size=12, color=WARN, align=PP_ALIGN.CENTER)
    add_text(slide, after, Inches(9.2), y + Inches(0.1), Inches(3.5), Inches(0.4),
             size=12, bold=True, color=color, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════
# 슬라이드 10: 다음 단계
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
header_bar(slide, "다음 단계: 기능 개선 10라운드", "동일 방법론으로 UX·기능·성능 개선 예정")
footer(slide, 10)

add_rect(slide, Inches(0.3), Inches(1.25), prs.slide_width - Inches(0.6), Inches(5.8),
         fill=LIGHT)

next_steps = [
    ("📊 대시보드 강화", "월별 목표 설정 & 달성률, KPI 카드 추가"),
    ("📅 세금 캘린더", "신고 일정 알림, 준비 체크리스트"),
    ("📈 고급 분석", "전년 동기 비교, 계절성 분석 차트"),
    ("🏷️ 커스텀 카테고리", "사용자 정의 카테고리 추가/수정"),
    ("📱 모바일 최적화", "터치 인터랙션, 반응형 레이아웃 개선"),
    ("🔔 알림 시스템", "세금 마감일 브라우저 알림"),
]

for i, (title, desc) in enumerate(next_steps):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.2)
    y = Inches(1.4 + row * 1.7)
    add_rect(slide, x, y, Inches(5.9), Inches(1.55), fill=WHITE)
    add_text(slide, title, x + Inches(0.15), y + Inches(0.12), Inches(5.6), Inches(0.45),
             size=14, bold=True, color=DARK)
    add_text(slide, desc, x + Inches(0.15), y + Inches(0.62), Inches(5.6), Inches(0.75),
             size=11, color=MID)

add_text(slide, "→ Gemini 리뷰 + ralph 루프로 10라운드 반복 예정",
         Inches(0.5), Inches(6.6), Inches(12), Inches(0.4),
         size=13, bold=True, color=BRAND)

# ── 저장 ──────────────────────────────────────────────
output = "/Users/taelkim/Jobs/boss-accounting/사장님회계_코드개선보고서.pptx"
prs.save(output)
print(f"✅ PPT 생성 완료: {output}")
