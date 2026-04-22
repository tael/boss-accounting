from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BRAND  = RGBColor(0x10, 0x99, 0x48)
DARK   = RGBColor(0x1F, 0x29, 0x37)
MID    = RGBColor(0x6B, 0x72, 0x80)
LIGHT  = RGBColor(0xF9, 0xFA, 0xFB)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLUE   = RGBColor(0x25, 0x63, 0xEB)
PURPLE = RGBColor(0x70, 0x55, 0xE3)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
RED    = RGBColor(0xDC, 0x26, 0x26)
TEAL   = RGBColor(0x06, 0xB6, 0xD4)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def rect(slide, l, t, w, h, fill=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=14, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def hbar(slide, title, sub=""):
    rect(slide, 0, 0, prs.slide_width, Inches(1.05), BRAND)
    txt(slide, title, Inches(0.4), Inches(0.1), Inches(11), Inches(0.55),
        size=26, bold=True, color=WHITE)
    if sub:
        txt(slide, sub, Inches(0.4), Inches(0.62), Inches(12), Inches(0.38),
            size=13, color=RGBColor(0xBB, 0xF7, 0xD0))

def foot(slide, n):
    rect(slide, 0, prs.slide_height - Inches(0.3), prs.slide_width, Inches(0.3),
         RGBColor(0xF3, 0xF4, 0xF6))
    txt(slide, "사장님 회계 도우미  |  소규모 사업자를 위한 로컬 회계 대시보드",
        Inches(0.4), prs.slide_height - Inches(0.28), Inches(11), Inches(0.26),
        size=9, color=MID)
    txt(slide, str(n), Inches(12.8), prs.slide_height - Inches(0.28),
        Inches(0.4), Inches(0.26), size=9, color=MID, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════
# 1. 표지
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
rect(s, 0, 0, prs.slide_width, prs.slide_height, DARK)
rect(s, 0, 0, Inches(0.4), prs.slide_height, BRAND)
rect(s, 0, prs.slide_height - Inches(2.2), prs.slide_width, Inches(2.2),
     RGBColor(0x11, 0x18, 0x27))

txt(s, "사장님 회계 도우미", Inches(0.8), Inches(1.2), Inches(12), Inches(1.1),
    size=48, bold=True, color=WHITE)
txt(s, "소규모 사업자를 위한 로컬 전용 회계 대시보드",
    Inches(0.8), Inches(2.4), Inches(11), Inches(0.55),
    size=22, color=RGBColor(0xBB, 0xF7, 0xD0))
txt(s, '"사장님이 알면 돈버는 회계" 책 기반 | 데이터 100% 로컬 저장',
    Inches(0.8), Inches(3.05), Inches(10), Inches(0.4),
    size=14, color=MID, italic=True)

kpis = [("6개 메뉴", "대시보드부터 세금까지"), ("로컬 전용", "서버 전송 없음"), ("2025 세율", "종합소득세·부가세")]
for i, (num, label) in enumerate(kpis):
    x = Inches(0.8 + i * 3.0)
    rect(s, x, Inches(4.0), Inches(2.7), Inches(1.3), RGBColor(0x37, 0x41, 0x51))
    txt(s, num, x, Inches(4.1), Inches(2.7), Inches(0.55),
        size=26, bold=True, color=BRAND, align=PP_ALIGN.CENTER)
    txt(s, label, x, Inches(4.7), Inches(2.7), Inches(0.4),
        size=11, color=MID, align=PP_ALIGN.CENTER)

txt(s, "React 19  ·  TypeScript  ·  Vite 6  ·  Zustand 5  ·  Tailwind CSS 4  ·  Recharts",
    Inches(0.8), Inches(5.6), Inches(12), Inches(0.35),
    size=11, color=RGBColor(0x4B, 0x55, 0x63))
txt(s, "2026.04.22", Inches(0.8), Inches(6.3), Inches(4), Inches(0.35),
    size=11, color=MID)

# ═══════════════════════════════════════════════
# 2. 기획 배경
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
hbar(s, "기획 배경", "소규모 사업자가 겪는 회계의 어려움을 해소한다")
foot(s, 2)

rect(s, 0, Inches(1.05), prs.slide_width, prs.slide_height - Inches(1.35), LIGHT)

problems = [
    ("😰 회계 용어가 어렵다", "매출·비용·손익을 직관적으로 파악하기 어려움"),
    ("📋 세금 신고가 막막하다", "부가세·종합소득세 언제 얼마 내야 하는지 모름"),
    ("💸 비용 관리가 안 된다", "어디에 얼마 쓰는지 한눈에 보이지 않음"),
    ("🔒 데이터 보안이 걱정된다", "재무 정보를 외부 서버에 올리기 꺼려짐"),
]
for i, (title, desc) in enumerate(problems):
    x = Inches(0.4 + (i % 2) * 6.2)
    y = Inches(1.25 + (i // 2) * 1.9)
    rect(s, x, y, Inches(5.9), Inches(1.7), WHITE)
    rect(s, x, y, Inches(0.15), Inches(1.7), RED)
    txt(s, title, x + Inches(0.3), y + Inches(0.15), Inches(5.4), Inches(0.45),
        size=14, bold=True, color=DARK)
    txt(s, desc, x + Inches(0.3), y + Inches(0.68), Inches(5.4), Inches(0.8),
        size=11, color=MID)

txt(s, "→  사장님 회계 도우미는 이 4가지 문제를 모두 해결합니다",
    Inches(0.5), Inches(5.2), Inches(12), Inches(0.45),
    size=14, bold=True, color=BRAND)

# ═══════════════════════════════════════════════
# 3. 핵심 기능 개요
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
hbar(s, "핵심 기능 6가지", "거래 기록부터 세금 시뮬레이터까지 All-in-One")
foot(s, 3)

features = [
    ("📊", "대시보드", BRAND,
     "이번 달 매출·비용·순이익 요약\n전월 대비 변화율 & 인사이트\n다가오는 세금 일정 알림"),
    ("📝", "거래 기록", BLUE,
     "매출/비용 CRUD (부가세 공제 여부 포함)\n카테고리 분류 & 검색/필터\n금액 입력 시 천단위 자동 포맷"),
    ("📈", "재무제표", PURPLE,
     "월별·분기별·연간 손익계산서\n현금흐름 차트 (최근 12개월)\n자산·부채 스냅샷"),
    ("🔍", "분석 도구", AMBER,
     "손익분기점(BEP) 계산기\n비용 구조 파이차트\n매출·비용 월별 추이 차트"),
    ("🧮", "세금 시뮬레이터", TEAL,
     "부가세 계산 (일반·간이 과세자)\n2025년 종합소득세 누진세율\n노란우산공제 절세 효과"),
    ("💾", "데이터 관리", RGBColor(0x84, 0xCC, 0x16),
     "JSON export/import 백업·복원\nGoogle Drive appData 동기화\n예시 데이터 38건 자동 생성"),
]

for i, (icon, name, color, detail) in enumerate(features):
    col, row = i % 3, i // 3
    x = Inches(0.3 + col * 4.3)
    y = Inches(1.2 + row * 2.9)
    rect(s, x, y, Inches(4.1), Inches(2.7), LIGHT)
    rect(s, x, y, Inches(4.1), Inches(0.55), color)
    txt(s, f"{icon}  {name}", x + Inches(0.15), y + Inches(0.08),
        Inches(3.8), Inches(0.42), size=15, bold=True, color=WHITE)
    txt(s, detail, x + Inches(0.15), y + Inches(0.65), Inches(3.8), Inches(1.9),
        size=10.5, color=DARK)

# ═══════════════════════════════════════════════
# 4. 거래 기록 & 대시보드
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
hbar(s, "거래 기록 & 대시보드", "직관적인 입력과 실시간 요약")
foot(s, 4)

# 왼쪽: 거래 입력 흐름
rect(s, Inches(0.3), Inches(1.15), Inches(5.9), Inches(5.95), LIGHT)
txt(s, "거래 입력 흐름", Inches(0.5), Inches(1.25), Inches(5.5), Inches(0.4),
    size=14, bold=True, color=DARK)

steps = [
    (BRAND, "① 유형 선택", "매출 / 비용 토글"),
    (BLUE, "② 날짜 입력", "오늘 날짜 자동 기본값 (로컬 기준)"),
    (PURPLE, "③ 금액 입력", "천단위 자동 포맷 (1,000,000원)"),
    (AMBER, "④ 카테고리", "용역매출·인건비 등 11개 분류"),
    (TEAL, "⑤ 저장", "localStorage 즉시 반영"),
]
for i, (color, step, desc) in enumerate(steps):
    y = Inches(1.75 + i * 0.95)
    rect(s, Inches(0.5), y, Inches(0.55), Inches(0.75), color)
    txt(s, step.split(" ")[0], Inches(0.5), y + Inches(0.1), Inches(0.55), Inches(0.55),
        size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, step[2:], Inches(1.15), y + Inches(0.05), Inches(2.2), Inches(0.35),
        size=12, bold=True, color=DARK)
    txt(s, desc, Inches(1.15), y + Inches(0.4), Inches(4.8), Inches(0.3),
        size=10, color=MID)

# 오른쪽: 대시보드 구성
rect(s, Inches(6.5), Inches(1.15), Inches(6.5), Inches(5.95), LIGHT)
txt(s, "대시보드 구성", Inches(6.7), Inches(1.25), Inches(6.2), Inches(0.4),
    size=14, bold=True, color=DARK)

dash_items = [
    (BRAND, "이번 달 요약 카드 4개", "매출 / 비용 / 순이익 / 이익률\n전월 대비 ▲▼ 변화율 표시"),
    (BLUE, "인사이트 카드", "인건비 비율 경고\n이익률 달성 축하 등 자동 분석"),
    (AMBER, "세금 일정", "이번 달 / 다음 달 신고 일정\n부가세·종소세·원천세 구분"),
    (TEAL, "최근 거래 5건", "날짜·메모·금액 빠른 조회"),
]
for i, (color, title, desc) in enumerate(dash_items):
    y = Inches(1.75 + i * 1.25)
    rect(s, Inches(6.7), y, Inches(0.18), Inches(1.1), color)
    txt(s, title, Inches(7.0), y + Inches(0.08), Inches(5.7), Inches(0.35),
        size=12, bold=True, color=DARK)
    txt(s, desc, Inches(7.0), y + Inches(0.48), Inches(5.7), Inches(0.55),
        size=10, color=MID)

# ═══════════════════════════════════════════════
# 5. 재무제표 & 분석
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
hbar(s, "재무제표 & 분석 도구", "손익계산서부터 손익분기점까지")
foot(s, 5)

sections = [
    ("📈 손익계산서", BRAND, [
        "기간: 월별 / 분기별 / 연간 선택",
        "총 매출 / 총 비용 / 순이익 계산",
        "카테고리별 비용 내역 상세 표시",
        "이익률 % 자동 산출",
    ]),
    ("💧 현금흐름 차트", BLUE, [
        "최근 12개월 월별 유입·유출 막대차트",
        "순 현금흐름 라인 오버레이",
        "Recharts 기반 반응형 렌더링",
        "Y축 단위 자동 조정 (만·억)",
    ]),
    ("📊 비용 구조 분석", PURPLE, [
        "비용 카테고리별 파이차트",
        "고정비 vs 변동비 분류 집계",
        "최근 3개월 평균 기준",
    ]),
    ("⚖️ 손익분기점(BEP)", AMBER, [
        "고정비 / 변동비 / 매출 입력",
        "BEP 매출액 자동 계산",
        "달성 불가 시 경고 표시",
        "계산식: 고정비 ÷ (1 - 변동비율)",
    ]),
]

for i, (title, color, items) in enumerate(sections):
    col, row = i % 2, i // 2
    x = Inches(0.3 + col * 6.4)
    y = Inches(1.2 + row * 2.85)
    rect(s, x, y, Inches(6.1), Inches(2.65), LIGHT)
    rect(s, x, y, Inches(6.1), Inches(0.45), color)
    txt(s, title, x + Inches(0.15), y + Inches(0.05), Inches(5.8), Inches(0.38),
        size=14, bold=True, color=WHITE)
    for j, item in enumerate(items):
        txt(s, f"• {item}", x + Inches(0.2), y + Inches(0.55 + j * 0.48),
            Inches(5.7), Inches(0.42), size=10.5, color=DARK)

# ═══════════════════════════════════════════════
# 6. 세금 시뮬레이터
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
hbar(s, "세금 시뮬레이터", "2025년 세율 기준 참고용 계산 (세무사 상담 권장)")
foot(s, 6)

rect(s, 0, Inches(1.05), prs.slide_width, prs.slide_height - Inches(1.35), LIGHT)

# 부가세
rect(s, Inches(0.3), Inches(1.2), Inches(5.9), Inches(4.8), WHITE)
rect(s, Inches(0.3), Inches(1.2), Inches(5.9), Inches(0.45), TEAL)
txt(s, "🧮 부가세 계산기", Inches(0.5), Inches(1.25), Inches(5.5), Inches(0.38),
    size=15, bold=True, color=WHITE)

vat_rows = [
    ("과세 유형", "일반과세자 / 간이과세자"),
    ("일반과세", "매출세액(10%) - 매입세액(10%)"),
    ("간이과세", "매출액 × 업종별 부가가치율 × 10%"),
    ("업종 예시", "소매 15% / 음식점 10% / 서비스 30%"),
    ("결과", "납부액 / 환급액 자동 산출"),
]
for i, (label, val) in enumerate(vat_rows):
    y = Inches(1.8 + i * 0.6)
    rect(s, Inches(0.4), y, Inches(1.5), Inches(0.5), LIGHT)
    txt(s, label, Inches(0.45), y + Inches(0.08), Inches(1.4), Inches(0.38),
        size=10, bold=True, color=TEAL)
    txt(s, val, Inches(2.1), y + Inches(0.08), Inches(3.8), Inches(0.38),
        size=10, color=DARK)

# 종합소득세
rect(s, Inches(6.7), Inches(1.2), Inches(6.2), Inches(4.8), WHITE)
rect(s, Inches(6.7), Inches(1.2), Inches(6.2), Inches(0.45), PURPLE)
txt(s, "📋 종합소득세 계산기", Inches(6.9), Inches(1.25), Inches(5.8), Inches(0.38),
    size=15, bold=True, color=WHITE)

brackets = [
    ("1,400만 이하", "6%", "0"),
    ("5,000만 이하", "15%", "126만"),
    ("8,800만 이하", "24%", "576만"),
    ("1.5억 이하", "35%", "1,544만"),
    ("3억 이하", "38%", "1,994만"),
    ("5억 이하", "40%", "2,594만"),
    ("10억 이하", "42%", "3,594만"),
    ("10억 초과", "45%", "6,594만"),
]
txt(s, "2025년 누진세율 구간", Inches(6.9), Inches(1.75), Inches(5.8), Inches(0.35),
    size=11, bold=True, color=MID)
for i, (range_, rate, ded) in enumerate(brackets[:6]):
    y = Inches(2.15 + i * 0.52)
    c = PURPLE if i < 2 else (RGBColor(0x9B, 0x7E, 0xFF) if i < 4 else MID)
    txt(s, f"• {range_:12} → {rate}  (공제 {ded})", Inches(6.9), y,
        Inches(5.8), Inches(0.45), size=10, color=DARK)

txt(s, "공제 항목: 부양가족 × 150만 + 국민연금 + 노란우산공제(최대 500만)",
    Inches(6.9), Inches(5.4), Inches(5.8), Inches(0.45),
    size=9.5, color=MID, italic=True)

txt(s, "⚠️  본 계산은 2025년 세율 기준 참고용이며 실제 신고 시 세무사 상담을 권장합니다.",
    Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.38),
    size=10, color=RED, italic=True)

# ═══════════════════════════════════════════════
# 7. 데이터 관리 & 보안
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
hbar(s, "데이터 관리 & 보안", "100% 로컬 저장 — 서버로 전송되지 않습니다")
foot(s, 7)

# 로컬 저장 강조
rect(s, Inches(0.3), Inches(1.15), prs.slide_width - Inches(0.6), Inches(1.5),
     RGBColor(0xF0, 0xFD, 0xF4))
rect(s, Inches(0.3), Inches(1.15), Inches(0.2), Inches(1.5), BRAND)
txt(s, "🔒  모든 재무 데이터는 브라우저 localStorage에만 저장됩니다",
    Inches(0.7), Inches(1.3), Inches(12), Inches(0.45),
    size=16, bold=True, color=BRAND)
txt(s, "외부 서버 전송 없음 · 인터넷 연결 불필요 · 개인정보 유출 위험 없음",
    Inches(0.7), Inches(1.78), Inches(12), Inches(0.38),
    size=12, color=MID)

# 3개 열
cols_data = [
    ("📤 JSON 백업", BLUE, [
        "전체 거래 데이터 JSON 파일 다운로드",
        "날짜별 파일명 자동 생성",
        "사용자 설정(세금 유형 등) 포함",
        "정기적 백업 권장 (5MB 한도)",
    ]),
    ("📥 JSON 복원", PURPLE, [
        "백업 파일로 전체 데이터 복원",
        "손상된 항목 자동 필터링",
        "앱 식별자 검증으로 타 앱 파일 거부",
        "복원 시 확인 대화상자 표시",
    ]),
    ("☁️ Google Drive", AMBER, [
        "appData 폴더에 자동 저장",
        "기기 간 동기화 가능",
        "OAuth2 인증 (토큰 메모리에만 보관)",
        "백업·복원 버튼 1클릭",
    ]),
]

for i, (title, color, items) in enumerate(cols_data):
    x = Inches(0.3 + i * 4.3)
    y = Inches(2.85)
    rect(s, x, y, Inches(4.1), Inches(3.9), LIGHT)
    rect(s, x, y, Inches(4.1), Inches(0.45), color)
    txt(s, title, x + Inches(0.15), y + Inches(0.05), Inches(3.8), Inches(0.38),
        size=14, bold=True, color=WHITE)
    for j, item in enumerate(items):
        txt(s, f"• {item}", x + Inches(0.15), y + Inches(0.6 + j * 0.75),
            Inches(3.8), Inches(0.65), size=10.5, color=DARK)

# ═══════════════════════════════════════════════
# 8. 기술 스택
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
hbar(s, "기술 스택", "최신 React 생태계 기반의 경량 로컬 앱")
foot(s, 8)

stacks = [
    ("Frontend", BRAND, [
        ("React 19", "Concurrent Features, Suspense"),
        ("TypeScript", "엄격한 타입 안전성"),
        ("Vite 6", "빠른 HMR & 번들링"),
        ("Tailwind CSS 4", "CSS-first 설정"),
    ]),
    ("상태 관리", BLUE, [
        ("Zustand 5", "persist middleware + migrate"),
        ("localStorage", "브라우저 로컬 저장"),
        ("스키마 버전 관리", "마이그레이션 지원"),
        ("rehydration 검증", "손상 데이터 자동 필터"),
    ]),
    ("차트 & 라우팅", PURPLE, [
        ("Recharts", "반응형 SVG 차트"),
        ("react-router 7", "HashRouter (GitHub Pages)"),
        ("react-focus-lock", "모달 포커스 트랩"),
        ("Google Identity", "OAuth2 드라이브 연동"),
    ]),
    ("테스트 & 빌드", AMBER, [
        ("Vitest", "단위 테스트 96건"),
        ("Testing Library", "컴포넌트 테스트"),
        ("GitHub Actions", "CI/CD 자동 배포"),
        ("GitHub Pages", "정적 배포"),
    ]),
]

for i, (cat, color, items) in enumerate(stacks):
    x = Inches(0.3 + i * 3.15)
    rect(s, x, Inches(1.15), Inches(3.0), Inches(5.95), LIGHT)
    rect(s, x, Inches(1.15), Inches(3.0), Inches(0.45), color)
    txt(s, cat, x + Inches(0.1), Inches(1.2), Inches(2.8), Inches(0.38),
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, (name, desc) in enumerate(items):
        y = Inches(1.72 + j * 1.25)
        rect(s, x + Inches(0.1), y, Inches(2.8), Inches(1.1), WHITE)
        txt(s, name, x + Inches(0.2), y + Inches(0.08), Inches(2.6), Inches(0.38),
            size=12, bold=True, color=DARK)
        txt(s, desc, x + Inches(0.2), y + Inches(0.52), Inches(2.6), Inches(0.45),
            size=9.5, color=MID)

# ═══════════════════════════════════════════════
# 9. 화면 구성 & 사용 흐름
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
hbar(s, "화면 구성 & 사용 흐름", "6개 페이지 — 사이드바 네비게이션")
foot(s, 9)

rect(s, 0, Inches(1.05), prs.slide_width, prs.slide_height - Inches(1.35), LIGHT)

# 사이드바 시뮬레이션
rect(s, Inches(0.3), Inches(1.15), Inches(2.2), Inches(5.95),
     RGBColor(0x1F, 0x29, 0x37))
txt(s, "📱 사이드바", Inches(0.4), Inches(1.25), Inches(2.0), Inches(0.35),
    size=10, color=MID)

nav_items = [
    ("📊", "대시보드", True),
    ("📝", "거래 기록", False),
    ("📈", "재무제표", False),
    ("🔍", "분석", False),
    ("🧮", "세금", False),
    ("💾", "데이터", False),
]
for i, (icon, name, active) in enumerate(nav_items):
    y = Inches(1.7 + i * 0.78)
    bg = BRAND if active else None
    if bg:
        rect(s, Inches(0.35), y, Inches(2.1), Inches(0.65), bg)
    c = WHITE if active else MID
    txt(s, f"{icon}  {name}", Inches(0.5), y + Inches(0.12), Inches(1.85), Inches(0.42),
        size=12, bold=active, color=c)

# 페이지 설명
pages = [
    ("대시보드", BRAND, "월 요약 4카드 + 인사이트 + 세금 일정 + 최근 거래"),
    ("거래 기록", BLUE, "인라인 폼 입력 + 필터 사이드바 + 정렬 목록"),
    ("재무제표", PURPLE, "손익계산서 + 현금흐름 차트 + 자산·부채 스냅샷"),
    ("분석", AMBER, "BEP 계산기 + 비용 파이차트 + 월별 추이"),
    ("세금", TEAL, "부가세 + 종합소득세 + 세금 팁 + 면책 고지"),
    ("데이터", RGBColor(0x84, 0xCC, 0x16), "예시데이터 생성 + JSON 백업/복원 + 구글드라이브"),
]
for i, (name, color, desc) in enumerate(pages):
    col, row = i % 2, i // 2
    x = Inches(2.8 + col * 5.1)
    y = Inches(1.2 + row * 1.8)
    rect(s, x, y, Inches(4.9), Inches(1.65), WHITE)
    rect(s, x, y, Inches(0.18), Inches(1.65), color)
    txt(s, name, x + Inches(0.3), y + Inches(0.12), Inches(4.4), Inches(0.38),
        size=13, bold=True, color=DARK)
    txt(s, desc, x + Inches(0.3), y + Inches(0.6), Inches(4.4), Inches(0.85),
        size=10, color=MID)

# ═══════════════════════════════════════════════
# 10. 로드맵
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
hbar(s, "향후 로드맵", "기능 개선 10라운드 — Gemini 리뷰 기반 반복 개선")
foot(s, 10)

rect(s, 0, Inches(1.05), prs.slide_width, prs.slide_height - Inches(1.35), LIGHT)

roadmap = [
    ("🎯 P0 — 핵심 UX", BRAND, [
        "월별 목표 매출 설정 & 달성률 표시",
        "거래 수정/삭제 인라인 편집",
        "모바일 터치 최적화",
    ]),
    ("📊 P1 — 분석 강화", BLUE, [
        "전년 동기 비교 차트",
        "카테고리별 예산 설정",
        "계절성 & 트렌드 분석",
    ]),
    ("🏷️ P2 — 커스터마이징", PURPLE, [
        "사용자 정의 카테고리 추가",
        "사업자 프로필 & 인보이스 생성",
        "다중 사업 계정 관리",
    ]),
    ("🔔 P3 — 알림 & 연동", AMBER, [
        "세금 마감일 브라우저 알림",
        "영수증 사진 첨부 (OCR)",
        "공공 API 세금 일정 연동",
    ]),
]

for i, (title, color, items) in enumerate(roadmap):
    col, row = i % 2, i // 2
    x = Inches(0.4 + col * 6.2)
    y = Inches(1.3 + row * 2.6)
    rect(s, x, y, Inches(5.9), Inches(2.4), WHITE)
    rect(s, x, y, Inches(5.9), Inches(0.45), color)
    txt(s, title, x + Inches(0.15), y + Inches(0.05), Inches(5.6), Inches(0.38),
        size=14, bold=True, color=WHITE)
    for j, item in enumerate(items):
        txt(s, f"→ {item}", x + Inches(0.2), y + Inches(0.6 + j * 0.58),
            Inches(5.5), Inches(0.5), size=11, color=DARK)

txt(s, "현재 진행 중: 클린코드 개선 8라운드 완료 (37개 스토리) → 기능 개선 10라운드 예정",
    Inches(0.4), Inches(6.65), Inches(12.5), Inches(0.38),
    size=11, bold=True, color=BRAND)

# ═══════════════════════════════════════════════
# 저장
# ═══════════════════════════════════════════════
out = "/Users/taelkim/Jobs/boss-accounting/사장님회계_프로젝트소개.pptx"
prs.save(out)
print(f"✅  PPT 생성 완료: {out}")
